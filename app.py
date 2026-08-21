import streamlit as st
import sqlite3
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from foundry_local_sdk import Configuration, FoundryLocalManager

from retriever import search
from chunker import chunk_text
from vector_store import create_database, save_embedding


# =========================================================
# PAGE CONFIGURATION
# =========================================================

st.set_page_config(
    page_title="Local RAG Assistant",
    page_icon="🤖",
    layout="wide"
)


# =========================================================
# DATABASE
# =========================================================

DATABASE_PATH = "rag_database.db"

create_database()


# =========================================================
# CUSTOM CSS
# =========================================================

st.markdown(
    """
    <style>

    /* =====================================================
       MAIN TITLE
       ===================================================== */

    .main-title {
        font-size: 38px;
        font-weight: 700;
        color: #1f2937;
        margin-bottom: 4px;
    }

    .subtitle {
        font-size: 16px;
        color: #6b7280;
        margin-bottom: 18px;
    }


    /* =====================================================
       LOCAL STATUS
       ===================================================== */

    .status {
        padding: 7px 13px;
        border-radius: 20px;
        background-color: #dcfce7;
        color: #166534;
        font-size: 13px;
        font-weight: 600;
        display: inline-block;
    }


    /* =====================================================
       SOURCE BOX
       ===================================================== */

    .source-box {
        padding: 12px 14px;
        border-radius: 10px;
        background-color: #f5f5f5;
        border: 1px solid #e5e7eb;
        margin-top: 8px;
        color: #374151;
    }


    /* =====================================================
       WARNING BOX
       ===================================================== */

    .warning-box {
        background-color: #fef2f2;
        border: 1px solid #fecaca;
        color: #991b1b;
        padding: 15px;
        border-radius: 12px;
        margin: 10px 0;
        font-weight: 500;
    }


    /* =====================================================
       CHAT INPUT
       ===================================================== */

    div[data-testid="stChatInput"] {
        border: 1px solid #e2e2e2 !important;
        border-radius: 16px !important;
        background-color: #ffffff !important;
        box-shadow: none !important;
    }

    div[data-testid="stChatInput"]:focus-within {
        border: 1px solid #d5d5d5 !important;
        box-shadow: 0 0 0 1px #eeeeee !important;
    }

    div[data-testid="stChatInput"] textarea {
        color: #222222 !important;
        background-color: transparent !important;
    }

    div[data-testid="stChatInput"] textarea::placeholder {
        color: #777777 !important;
        opacity: 1 !important;
    }

    div[data-testid="stChatInput"] button {
        color: #555555 !important;
        background-color: transparent !important;
    }


    /* =====================================================
       BUTTONS
       ===================================================== */

    div.stButton > button {
        border-radius: 9px !important;
        border: 1px solid #d1d5db !important;
        font-weight: 600 !important;
    }

    div.stButton > button:hover {
        border-color: #9ca3af !important;
        color: #374151 !important;
    }


    /* =====================================================
       FOOTER
       ===================================================== */

    .footer {
        text-align: center;
        color: #9ca3af;
        font-size: 12px;
        margin-top: 30px;
        padding-top: 15px;
        border-top: 1px solid #e5e7eb;
        line-height: 1.7;
    }

    </style>
    """,
    unsafe_allow_html=True
)


# =========================================================
# FOUNDRY LOCAL INITIALIZATION
# =========================================================

@st.cache_resource
def initialize_models():

    config = Configuration(
        app_name="rag-assistant",
        model_cache_dir=r"C:\Users\nurll\.foundry\cache\models"
    )

    manager = FoundryLocalManager.instance

    if manager is None:

        try:
            FoundryLocalManager.initialize(config)
        except Exception:
            pass

        manager = FoundryLocalManager.instance

    if manager is None:
        raise RuntimeError(
            "Foundry Local Manager could not be initialized."
        )

    # -----------------------------------------------------
    # Embedding model
    # -----------------------------------------------------

    embedding_model = manager.catalog.get_model(
        "qwen3-embedding-0.6b"
    )

    embedding_model.load()

    embedding_client = (
        embedding_model.get_embedding_client()
    )

    # -----------------------------------------------------
    # Generation model
    # -----------------------------------------------------

    generation_model = manager.catalog.get_model(
        "phi-3.5-mini"
    )

    generation_model.load()

    chat_client = (
        generation_model.get_chat_client()
    )

    return embedding_client, chat_client


# =========================================================
# LOAD MODELS
# =========================================================

with st.spinner("Loading local AI models..."):

    embedding_client, chat_client = initialize_models()


# =========================================================
# HEADER
# =========================================================

st.markdown(
    '<div class="main-title">🤖 Local RAG Assistant</div>',
    unsafe_allow_html=True
)

st.markdown(
    '<div class="subtitle">'
    'A private, document-based AI assistant powered by Microsoft Foundry Local'
    '</div>',
    unsafe_allow_html=True
)

st.markdown(
    '<span class="status">● LOCAL / OFFLINE</span>',
    unsafe_allow_html=True
)

st.divider()


# =========================================================
# SIDEBAR
# =========================================================

with st.sidebar:

    st.header("⚙️ System")

    st.markdown("### Generation Model")
    st.info("Phi-3.5-mini")

    st.markdown("### Embedding Model")
    st.info("Qwen3-Embedding-0.6B")

    st.markdown("### Knowledge Base")

    connection = sqlite3.connect(
        DATABASE_PATH
    )

    cursor = connection.cursor()

    cursor.execute(
        "SELECT COUNT(DISTINCT source) FROM documents"
    )

    document_count = cursor.fetchone()[0]

    cursor.execute(
        "SELECT COUNT(*) FROM documents"
    )

    chunk_count = cursor.fetchone()[0]

    cursor.execute(
        "SELECT DISTINCT source FROM documents"
    )

    database_documents = [
        row[0]
        for row in cursor.fetchall()
    ]

    connection.close()

    st.metric(
        "Documents",
        document_count
    )

    st.metric(
        "Chunks",
        chunk_count
    )

    st.divider()


    # =====================================================
    # DOCUMENT UPLOAD
    # =====================================================

    st.markdown("### 📄 Add Document")

    uploaded_file = st.file_uploader(
        "Upload a document",
        type=[
            "txt",
            "pdf",
            "docx",
            "pptx"
        ],
        help="Supported formats: TXT, PDF, DOCX and PPTX"
    )

    if uploaded_file is not None:

        st.caption(
            f"Selected: {uploaded_file.name}"
        )

        if st.button(
            "Process Document",
            use_container_width=True
        ):

            document_name = uploaded_file.name

            file_extension = (
                Path(document_name)
                .suffix
                .lower()
            )

            document_path = (
                Path("documents")
                / document_name
            )

            document_path.parent.mkdir(
                exist_ok=True
            )

            file_text = ""


            # =================================================
            # TXT
            # =================================================

            if file_extension == ".txt":

                try:

                    file_text = (
                        uploaded_file
                        .read()
                        .decode("utf-8")
                    )

                except UnicodeDecodeError:

                    st.error(
                        "The TXT file could not be read as UTF-8."
                    )

                    st.stop()


            # =================================================
            # PDF
            # =================================================

            elif file_extension == ".pdf":

                try:

                    pdf_reader = PdfReader(
                        uploaded_file
                    )

                    pages = []

                    for page in pdf_reader.pages:

                        text = page.extract_text()

                        if text:
                            pages.append(text)

                    file_text = (
                        "\n\n".join(pages)
                    )

                except Exception as error:

                    st.error(
                        f"PDF could not be processed: {error}"
                    )

                    st.stop()


            # =================================================
            # DOCX
            # =================================================

            elif file_extension == ".docx":

                try:

                    document = Document(
                        uploaded_file
                    )

                    paragraphs = []

                    for paragraph in document.paragraphs:

                        if paragraph.text.strip():

                            paragraphs.append(
                                paragraph.text
                            )

                    file_text = (
                        "\n\n".join(paragraphs)
                    )

                except Exception as error:

                    st.error(
                        f"Word document could not be processed: {error}"
                    )

                    st.stop()


            # =================================================
            # PPTX
            # =================================================

            elif file_extension == ".pptx":

                try:

                    presentation = Presentation(
                        uploaded_file
                    )

                    slides_text = []

                    for slide in presentation.slides:

                        slide_parts = []

                        for shape in slide.shapes:

                            if hasattr(shape, "text"):

                                text = shape.text.strip()

                                if text:

                                    slide_parts.append(
                                        text
                                    )

                        if slide_parts:

                            slides_text.append(
                                "\n".join(slide_parts)
                            )

                    file_text = (
                        "\n\n".join(slides_text)
                    )

                except Exception as error:

                    st.error(
                        f"PowerPoint presentation could not be processed: {error}"
                    )

                    st.stop()


            # =================================================
            # CHECK TEXT
            # =================================================

            if not file_text.strip():

                st.error(
                    "No readable text was found in this document."
                )

            else:

                with st.spinner(
                    "Processing document..."
                ):

                    # -------------------------------------------------
                    # Save original document
                    # -------------------------------------------------

                    document_path.write_bytes(
                        uploaded_file.getvalue()
                    )


                    # -------------------------------------------------
                    # Remove old version if it exists
                    # -------------------------------------------------

                    connection = sqlite3.connect(
                        DATABASE_PATH
                    )

                    cursor = connection.cursor()

                    cursor.execute(
                        """
                        DELETE FROM documents
                        WHERE source = ?
                        """,
                        (
                            str(document_path),
                        )
                    )

                    connection.commit()
                    connection.close()


                    # -------------------------------------------------
                    # Create chunks
                    # -------------------------------------------------

                    chunks = chunk_text(
                        file_text
                    )


                    # -------------------------------------------------
                    # Generate embeddings
                    # -------------------------------------------------

                    for index, chunk in enumerate(
                        chunks,
                        start=1
                    ):

                        response = (
                            embedding_client
                            .generate_embedding(
                                chunk
                            )
                        )

                        vector = (
                            response
                            .data[0]
                            .embedding
                        )

                        save_embedding(
                            source=str(
                                document_path
                            ),
                            chunk_id=index,
                            content=chunk,
                            embedding=vector
                        )


                st.success(
                    f"'{document_name}' processed successfully!"
                )

                st.info(
                    f"{len(chunks)} chunks added to the knowledge base."
                )

                st.rerun()


    # =====================================================
    # DOCUMENT LIST
    # =====================================================

    st.markdown("### 📚 Documents")

    if database_documents:

        for document in database_documents:

            document_path = Path(
                document
            )

            document_name = (
                document_path.name
            )

            extension = (
                document_path.suffix
                .lower()
            )


            # -------------------------------------------------
            # File icons
            # -------------------------------------------------

            if extension == ".pdf":

                icon = "📖"

            elif extension == ".docx":

                icon = "📘"

            elif extension == ".pptx":

                icon = "📙"

            elif extension == ".txt":

                icon = "🗒️"

            else:

                icon = "📄"


            # -------------------------------------------------
            # Document row
            # -------------------------------------------------

            col1, col2 = st.columns(
                [5, 1]
            )

            with col1:

                st.write(
                    f"{icon} {document_name}"
                )

            with col2:

                if st.button(
                    "×",
                    key=f"delete_{document}",
                    help=f"Delete {document_name}"
                ):

                    # Delete from database

                    connection = sqlite3.connect(
                        DATABASE_PATH
                    )

                    cursor = connection.cursor()

                    cursor.execute(
                        """
                        DELETE FROM documents
                        WHERE source = ?
                        """,
                        (
                            document,
                        )
                    )

                    connection.commit()
                    connection.close()


                    # Delete physical document

                    try:

                        if document_path.exists():

                            document_path.unlink()

                    except Exception:
                        pass


                    st.rerun()

    else:

        st.caption(
            "No documents available."
        )


    st.divider()


    # =====================================================
    # CLEAR CHAT
    # =====================================================

    st.markdown("### 💬 Chat")

    if st.button(
        "🗑️ Clear Chat",
        use_container_width=True
    ):

        st.session_state.messages = []

        st.rerun()


    st.divider()


    # =====================================================
    # FOOTER
    # =====================================================

    st.markdown(
        '<div class="footer">'
        'Powered by Microsoft Foundry Local<br>'
        'All inference runs locally.<br><br>'
        '<b>Developed by N. Çiftci · 2026</b>'
        '</div>',
        unsafe_allow_html=True
    )


# =========================================================
# MAIN CHAT
# =========================================================

st.subheader(
    "💬 Ask your documents"
)

st.write(
    "Ask questions about the information available "
    "in your local knowledge base."
)


# =========================================================
# CHAT HISTORY
# =========================================================

if "messages" not in st.session_state:

    st.session_state.messages = []


for message in st.session_state.messages:

    with st.chat_message(
        message["role"]
    ):

        st.markdown(
            message["content"]
        )

        if message.get("sources"):

            st.markdown(
                "#### 📚 Sources"
            )

            displayed_sources = []

            for source in message["sources"]:

                source_name = Path(
                    source["source"]
                ).name

                if source_name not in displayed_sources:

                    displayed_sources.append(
                        source_name
                    )

                    st.markdown(
                        f"""
                        <div class="source-box">
                        📄 <b>{source_name}</b>
                        </div>
                        """,
                        unsafe_allow_html=True
                    )


# =========================================================
# QUESTION INPUT
# =========================================================

question = st.chat_input(
    "Ask a question about your documents..."
)


# =========================================================
# PROCESS QUESTION
# =========================================================

if question:

    # -----------------------------------------------------
    # USER MESSAGE
    # -----------------------------------------------------

    st.session_state.messages.append(
        {
            "role": "user",
            "content": question
        }
    )

    with st.chat_message("user"):

        st.markdown(
            question
        )


    # -----------------------------------------------------
    # ASSISTANT
    # -----------------------------------------------------

    with st.chat_message("assistant"):

        # =================================================
        # RETRIEVAL
        # =================================================

        with st.spinner(
            "🔎 Retrieving relevant information..."
        ):

            results = search(
                question,
                embedding_client,
                top_k=3,
                similarity_threshold=0.55
            )


        # =================================================
        # NO RESULTS
        # =================================================

        if not results:

            answer = (
                "I couldn't find relevant information "
                "in the local knowledge base."
            )

            st.markdown(
                f"""
                <div class="warning-box">
                ⚠️ {answer}
                </div>
                """,
                unsafe_allow_html=True
            )

            st.session_state.messages.append(
                {
                    "role": "assistant",
                    "content": answer,
                    "sources": []
                }
            )

            st.stop()


        # =================================================
        # CREATE CONTEXT
        # =================================================

        context_parts = []

        for result in results:

            context_parts.append(
                f"Source: {result['source']}\n"
                f"Chunk: {result['chunk_id']}\n"
                f"Similarity: {result['score']:.4f}\n"
                f"Content:\n{result['content']}"
            )

        context = (
            "\n\n".join(
                context_parts
            )
        )


        # =================================================
        # RAG PROMPT
        # =================================================

        messages = [

            {
                "role": "system",
                "content": (
                    "You are a helpful document assistant. "
                    "Answer the user's question using ONLY "
                    "the information provided in the context. "
                    "Do not use outside knowledge. "
                    "Do not invent facts. "
                    "If the context does not contain enough "
                    "information, say that the information "
                    "could not be found in the local "
                    "knowledge base."
                )
            },

            {
                "role": "user",
                "content": (
                    f"Context:\n\n"
                    f"{context}\n\n"
                    f"Question: {question}"
                )
            }

        ]


        # =================================================
        # GENERATE ANSWER
        # =================================================

        with st.spinner(
            "🤖 Generating answer..."
        ):

            response = (
                chat_client
                .complete_chat(
                    messages
                )
            )

        answer = (
            response
            .choices[0]
            .message
            .content
        )


        # =================================================
        # DISPLAY ANSWER
        # =================================================

        st.markdown(
            answer
        )


        # =================================================
        # SOURCES
        # =================================================

        st.markdown(
            "#### 📚 Sources"
        )

        displayed_sources = []

        for result in results:

            source_name = Path(
                result["source"]
            ).name

            if source_name not in displayed_sources:

                displayed_sources.append(
                    source_name
                )

                st.markdown(
                    f"""
                    <div class="source-box">
                    📄 <b>{source_name}</b>
                    </div>
                    """,
                    unsafe_allow_html=True
                )


        # =================================================
        # SAVE ASSISTANT MESSAGE
        # =================================================

        st.session_state.messages.append(
            {
                "role": "assistant",
                "content": answer,
                "sources": results
            }
        )